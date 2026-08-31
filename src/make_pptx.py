"""
Rapport PowerPoint du Défi 2 — 10 diapositives, générées depuis data/gold/.

    python src/make_pptx.py   ->  report/Defi2_Togo_Energie_Climat_Forets.pptx

Pourquoi un script plutôt qu'un fichier composé à la main : le deck et le
tableau de bord lisent **la même source**. Un chiffre ne peut donc pas
diverger entre l'écran et la diapositive, et une donnée corrigée dans
data/gold/ se propage aux deux d'une seule commande. Les couleurs viennent
de `app/palette.py`, pour la même raison.

L'ordre des diapositives est celui des pages du tableau de bord : un jury qui
tient le rapport retrouve l'écran page par page.
"""
import json
import sys
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_TICK_MARK
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
GOLD = ROOT / "data" / "gold"
OUT = ROOT / "report" / "Defi2_Togo_Energie_Climat_Forets.pptx"
# Armoiries en PNG transparent : PowerPoint affiche du SVG, mais
# python-pptx ne sait pas en embarquer. Voir src/make_armoiries_png.py.
ARMOIRIES = ROOT / "app" / "assets" / "armoiries_togo.png"

sys.path.insert(0, str(ROOT / "app"))
from palette import C, FAMILLE_BUREAU   # noqa: E402

FAM = FAMILLE_BUREAU
AUTEUR = "KOUTSAVA Kossi Epiphane"
SOURCE = ("Six jeux de données du défi · datalab.gouv.tg — Banque Mondiale, "
          "inventaire national des GES 2018, stations météorologiques, "
          "forêts classées")

nat = json.load(open(GOLD / "diagnostic_national.json", encoding="utf-8"))
R = nat["reperes"]
S = nat["series"]
er, ec, rs = R["elec_rural"], R["ecart_urbain_rural"], R["ruraux_sans_elec"]
cb, dfr, sa = R["combustibles"], R["deforestation"], R["sante"]
rp, fi, ges, fo = (R["renouvelable_piege"], R["fiabilite"], R["ges"],
                   R["forets"])
tr, rg = R["tapis_roulant"], R["regimes_foret"]
PRG = {"CO2": 1, "CH4": 28, "N2O": 265}          # PRG 100 ans, GIEC AR5

# La portée de chaque levier sert deux fois : en synthèse et en détail.
POP = sorted((int(a), v) for a, v in S["pop_totale"].items()
             if v is not None)[-1][1]
PERS_BIOMASSE = POP * cb["biomasse"][1] / 100

L, H = Inches(13.333), Inches(7.5)               # 16:9
MARGE = Inches(.62)
LARG = L - 2 * MARGE
# La zone de titre réserve deux lignes sur toutes les diapositives, et le
# contenu démarre toujours à la même hauteur : en feuilletant, le regard
# n'a pas à se recaler d'une diapositive à l'autre.
HAUT = Inches(1.72)


def fr(x, d=0):
    """Format français, espace insécable pour les milliers."""
    return f"{x:,.{d}f}".replace(",", " ").replace(".", ",")


def rgb(cle):
    h = C[cle].lstrip("#")
    return RGBColor(*(int(h[i:i + 2], 16) for i in (0, 2, 4)))


BLANC = RGBColor(0xFF, 0xFF, 0xFF)


# ------------------------------------------------------------ primitives
def texte(slide, x, y, w, h, contenu, taille=12, gras=False, couleur="encre",
          align=PP_ALIGN.LEFT, interligne=1.25, majuscules=False,
          espacement=None, ancre=MSO_ANCHOR.TOP):
    """Bloc de texte. `contenu` : une chaîne, ou une liste de paragraphes."""
    boite = slide.shapes.add_textbox(x, y, w, h)
    tf = boite.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = ancre
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    lignes = contenu if isinstance(contenu, list) else [contenu]
    for i, ligne in enumerate(lignes):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = interligne
        if i:
            p.space_before = Pt(6)
        # une ligne peut être (texte, couleur, gras) pour teinter un fragment
        morceaux = ligne if isinstance(ligne, list) else [ligne]
        for m in morceaux:
            t, coul, g = (m if isinstance(m, tuple) else (m, couleur, gras))
            r = p.add_run()
            r.text = t.upper() if majuscules else t
            r.font.size = Pt(taille)
            r.font.bold = g
            r.font.name = FAM
            r.font.color.rgb = rgb(coul) if isinstance(coul, str) else coul
            if espacement:
                _espacer(r, espacement)
    return boite


def _espacer(run, points):
    """Interlettrage — python-pptx ne l'expose pas, on l'écrit dans le XML."""
    run.font._rPr.set("spc", str(int(points * 100)))


def rectangle(slide, x, y, w, h, fond=None, bord=None, arrondi=True):
    forme = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if arrondi else MSO_SHAPE.RECTANGLE,
        x, y, w, h)
    if arrondi:
        forme.adjustments[0] = .05
    if fond:
        forme.fill.solid()
        forme.fill.fore_color.rgb = rgb(fond) if isinstance(fond, str) else fond
    else:
        forme.fill.background()
    if bord:
        forme.line.color.rgb = rgb(bord)
        forme.line.width = Pt(.75)
    else:
        forme.line.fill.background()
    forme.shadow.inherit = False          # aucune ombre : à-plats francs
    forme.text_frame.word_wrap = True
    return forme


# --------------------------------------------------------------- gabarit
def diapo(prs, kicker, titre, numero):
    """Diapositive de contenu : intitulé de rubrique, titre-trouvaille, pied."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fond = rectangle(s, 0, 0, L, H, fond="surface", arrondi=False)
    fond.line.fill.background()
    texte(s, MARGE, Inches(.46), LARG, Inches(.22), kicker, taille=9.5,
          gras=True, couleur="sourdine", majuscules=True, espacement=.7)
    texte(s, MARGE, Inches(.72), LARG, Inches(.92), titre, taille=25.5,
          gras=True, couleur="encre", interligne=1.06)
    # Le pied de page laisse sa place à l'emblème, posé juste à gauche.
    texte(s, MARGE + Inches(.36), H - Inches(.46), LARG, Inches(.2),
          "Défi 2 · Énergie, Climat & Forêts au Togo", taille=8.5,
          couleur="sourdine")
    texte(s, L - MARGE - Inches(.6), H - Inches(.46), Inches(.6), Inches(.2),
          f"{numero} / 10", taille=8.5, couleur="sourdine",
          align=PP_ALIGN.RIGHT)
    if ARMOIRIES.exists():
        s.shapes.add_picture(str(ARMOIRIES), MARGE - Inches(.02),
                             H - Inches(.58), height=Inches(.34))
    return s


def tuile(slide, x, y, w, label, valeur, unite, sous, couleur,
          h=Inches(1.62)):
    """Chiffre-clé — même anatomie que le tableau de bord : intitulé, nombre,
    contexte. La couleur ne se pose que sur le nombre."""
    rectangle(slide, x, y, w, h, fond="surface", bord="bord")
    pad = Inches(.17)
    texte(slide, x + pad, y + Inches(.15), w - 2 * pad, Inches(.2), label,
          taille=8, gras=True, couleur="encre", majuscules=True, espacement=.5)
    texte(slide, x + pad, y + Inches(.4), w - 2 * pad, Inches(.45),
          [[(valeur, couleur, True), (f" {unite}" if unite else "",
                                      couleur, True)]],
          taille=26, interligne=1)
    texte(slide, x + pad, y + Inches(.92), w - 2 * pad, h - Inches(1.02),
          sous, taille=8.5, couleur="sourdine", interligne=1.3)


ENCARTS = {"constat": ("foret", "foret_l", "Constat"),
           "alerte": ("risque", "risque_l", "Alerte"),
           "action": ("energie", "energie_l", "Décision")}


def encart(slide, x, y, w, h, genre, corps, titre=None):
    """Bloc de lecture teinté — même grammaire que le tableau de bord."""
    coul, fond, defaut = ENCARTS[genre]
    rectangle(slide, x, y, w, h, fond=fond, bord=None)
    pad = Inches(.2)
    texte(slide, x + pad, y + Inches(.15), w - 2 * pad, Inches(.18),
          titre or defaut, taille=8.5, gras=True, couleur=coul,
          majuscules=True, espacement=.6)
    texte(slide, x + pad, y + Inches(.4), w - 2 * pad, h - Inches(.5),
          corps, taille=11, couleur="encre", interligne=1.35)


# ---------------------------------------------------------------- graphes
def _habiller(chart, couleurs, epaisseur=Pt(2.4), legende=True,
              format_val='0"%"', maxi=None, mini=None, unite_val=None):
    chart.font.size = Pt(9.5)
    chart.font.name = FAM
    chart.font.color.rgb = rgb("encre_2")
    chart.has_title = False
    chart.has_legend = legende
    if legende:
        chart.legend.position = XL_LEGEND_POSITION.TOP
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(9.5)

    va = chart.value_axis
    va.has_major_gridlines = True
    va.major_gridlines.format.line.color.rgb = rgb("bord")
    va.major_gridlines.format.line.width = Pt(.5)
    va.format.line.fill.background()
    va.tick_labels.number_format = format_val
    va.tick_labels.number_format_is_linked = False
    va.major_tick_mark = XL_TICK_MARK.NONE
    va.minor_tick_mark = XL_TICK_MARK.NONE
    if maxi is not None:
        va.maximum_scale = maxi
    if mini is not None:
        va.minimum_scale = mini
    if unite_val:
        va.major_unit = unite_val

    ca = chart.category_axis
    ca.has_major_gridlines = False
    ca.format.line.color.rgb = rgb("bord_fort")
    ca.major_tick_mark = XL_TICK_MARK.NONE
    ca.minor_tick_mark = XL_TICK_MARK.NONE
    return chart


def courbes(slide, x, y, w, h, annees, series, maxi=None, mini=None,
            legende=True, format_val='0"%"', pas_etiquette=5):
    """Courbes annuelles. Une année sur `pas_etiquette` est étiquetée :
    au-delà, l'axe devient illisible et on ne lit plus rien."""
    cd = CategoryChartData()
    cd.categories = [str(a) if a % pas_etiquette == 0 else "" for a in annees]
    for nom, valeurs, _ in series:
        cd.add_series(nom, valeurs)
    graphe = slide.shapes.add_chart(XL_CHART_TYPE.LINE, x, y, w, h, cd).chart
    for serie, (_, _, coul) in zip(graphe.series, series):
        serie.smooth = False
        serie.format.line.color.rgb = rgb(coul)
        serie.format.line.width = Pt(2.4)
    return _habiller(graphe, [s[2] for s in series], legende=legende,
                     format_val=format_val, maxi=maxi, mini=mini)


def barres(slide, x, y, w, h, categories, series, horizontal=True,
           legende=True, format_val='0"%"', maxi=None):
    cd = CategoryChartData()
    cd.categories = categories
    for nom, valeurs, _ in series:
        cd.add_series(nom, valeurs)
    genre = (XL_CHART_TYPE.BAR_CLUSTERED if horizontal
             else XL_CHART_TYPE.COLUMN_CLUSTERED)
    graphe = slide.shapes.add_chart(genre, x, y, w, h, cd).chart
    graphe.gap_width = 60
    graphe.overlap = -12
    for serie, (_, _, coul) in zip(graphe.series, series):
        if isinstance(coul, list):                 # une teinte par barre
            for point, c in zip(serie.points, coul):
                point.format.fill.solid()
                point.format.fill.fore_color.rgb = rgb(c)
        else:
            serie.format.fill.solid()
            serie.format.fill.fore_color.rgb = rgb(coul)
        serie.format.line.fill.background()
    return _habiller(graphe, None, legende=legende, format_val=format_val,
                     maxi=maxi)


def serie_gold(cle, debut=None):
    """Série annuelle de data/gold, en (années, valeurs)."""
    points = sorted((int(a), v) for a, v in S[cle].items() if v is not None)
    if debut:
        points = [(a, v) for a, v in points if a >= debut]
    return [a for a, _ in points], [v for _, v in points]


def aligner(annees_ref, annees, valeurs):
    """Cale une série sur l'axe commun, trous compris."""
    d = dict(zip(annees, valeurs))
    return [d.get(a) for a in annees_ref]


# --------------------------------------------------- contenu partage
# Les cinq chiffres du diagnostic et les trois leviers servent deux
# fois chacun : en synthese (diapositive 2) et en detail (10). Une
# seule definition, donc aucun risque de divergence entre les deux.
CHIFFRES = [
    ("Accès rural à l'électricité", f"{er['valeur']:.0f}", "%",
     f"contre {ec['urbain']:.0f} % en ville en {ec['annee']} — "
     f"{ec['valeur']:.0f} points d'écart", "energie"),
    ("Ruraux privés d'électricité", f"+{tr['variation_pct']:.0f}", "%",
     f"depuis {tr['annee_debut']} — de {fr(tr['sans_elec_debut']/1e6, 2)} à "
     f"{fr(tr['sans_elec_fin']/1e6, 2)} M de personnes, alors que le taux "
     f"d'accès était multiplié par "
     f"{er['valeur']/er['valeur_depart']:.0f}", "risque"),
    ("Ménages au bois ou au charbon", f"{cb['biomasse'][1]:.0f}", "%",
     f"en {cb['annees'][1]}, contre {cb['biomasse'][0]:.0f} % en "
     f"{cb['annees'][0]} : aucun recul", "risque"),
    ("Forêt perdue chaque année", f"{fr(dfr['perte_actuelle_ha_an'])}", "ha",
     f"dans le régime en cours depuis {dfr['regime_depuis']} — trois fois moins "
     f"que dans les années 1990, mais sans une seule année de reprise",
     "risque"),
    ("Accélération requise", f"× {er['facteur_acceleration']:.0f}", "",
     f"il faudrait +{fr(er['rythme_requis_2030'], 1)} pt/an au lieu de "
     f"+{fr(er['rythme_observe'], 1)}", "risque"),
]
LEVIERS = [
    ("1", "risque", "Cuisson propre",
     "GPL subventionné en zone périurbaine, foyers améliorés certifiés en "
     "zone rurale enclavée.",
     f"{fr(PERS_BIOMASSE/1e6, 1)} M de personnes",
     f"De {fr(R['cuisson_rurale']['valeur'], 1)} % à 15 % d'accès rural à "
     f"une cuisson propre en 2030",
     "EG.CFT.ACCS.RU.ZS — et non le taux d'énergie renouvelable"),
    ("2", "energie", "Solaire villageois décentralisé",
     "Mini-réseaux avec stockage pour les gros bourgs, kits domestiques "
     "pour l'habitat dispersé, dimensionnés sur les usages productifs.",
     f"{fr(tr['sans_elec_fin']/1e6, 1)} M de ruraux sans électricité",
     f"Dépasser {fr(round(tr['seuil_stagnation'], -2))} raccordements ruraux "
     f"par an — le seuil démographique — puis viser l'accès universel",
     "EG.ELC.ACCS.RU.ZS + un indicateur de disponibilité horaire, "
     "absent des données actuelles"),
    ("3", "foret", "Protection ciblée des massifs",
     "Surveillance, régénération assistée et agroforesterie concentrées "
     "sur les massifs prioritaires plutôt que saupoudrées.",
     f"{fo['nb_robustes']} massifs · {fr(fo['surface_totale_ha'])} ha "
     "de forêts classées",
     "Perte nette nulle sur les massifs prioritaires",
     "AG.LND.FRST.K2 + un suivi surfacique par massif, à créer"),
]

# ============================================================ le deck
def construire():
    prs = Presentation()
    prs.slide_width, prs.slide_height = L, H

    # ---------------------------------------------------- 1. couverture
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rectangle(s, 0, 0, L, H, fond="nuit", arrondi=False)
    if ARMOIRIES.exists():
        s.shapes.add_picture(str(ARMOIRIES), L - Inches(3.6), Inches(1.5),
                             height=Inches(3.5))
    texte(s, MARGE, Inches(1.45), Inches(9.0), Inches(.24),
          "Défi 2 · datalab.gouv.tg · République togolaise", taille=11,
          gras=True, couleur="sur_nuit_2", majuscules=True, espacement=.9)
    texte(s, MARGE, Inches(1.95), Inches(8.8), Inches(1.8),
          "Électrifier les campagnes\nsans brûler les forêts", taille=42,
          gras=True, couleur=BLANC, interligne=1.06)
    texte(s, MARGE, Inches(3.85), Inches(8.5), Inches(.95),
          "Le Togo vise l'accès universel à l'électricité en 2030. Les données "
          "montrent que l'objectif ne se joue pas sur le réseau : il se joue "
          "sur la cuisson des ménages, qui consomme la forêt.",
          taille=14, couleur="sur_nuit", interligne=1.5)

    rep = [("Accès rural à l'électricité", f"{er['valeur']:.0f} %"),
           ("Ménages au bois ou au charbon", f"{cb['biomasse'][1]:.0f} %"),
           ("Forêt perdue chaque année",
            f"−{fr(dfr['perte_actuelle_ha_an'])} ha")]
    for i_r, (lab, val) in enumerate(rep):
        x = MARGE + i_r * Inches(3.0)
        texte(s, x, Inches(5.15), Inches(2.85), Inches(.2), lab, taille=9,
              gras=True, couleur="sur_nuit_2", majuscules=True, espacement=.6)
        texte(s, x, Inches(5.4), Inches(2.85), Inches(.5), val, taille=28,
              gras=True, couleur=BLANC, interligne=1)
    # Le nom porte le travail : il se lit sans qu'on ait à le chercher.
    texte(s, MARGE, H - Inches(1.38), Inches(8.5), Inches(.4), AUTEUR,
          taille=17, gras=True, couleur=BLANC)
    texte(s, MARGE, H - Inches(.9), Inches(9.3), Inches(.6), SOURCE,
          taille=9, couleur="sur_nuit_2", interligne=1.4)

    # ------------------------------------------- 2. ce qui a été analysé
    s = diapo(prs, "Méthode", "Six jeux de données, une seule question : "
              "où agir en premier ?", 2)
    y = HAUT
    jeux = [
        ("Accès à l'électricité", "Banque Mondiale · total, urbain, rural · "
         f"{er['annee_depart']}–{er['annee']}"),
        ("Cuisson des ménages", "Combustible principal (deux enquêtes "
         f"{cb['annees'][0]} et {cb['annees'][1]}) et accès à une cuisson "
         "propre, série annuelle"),
        ("Énergies renouvelables", "Part des renouvelables dans la "
         "consommation finale, croisée avec la cuisson propre"),
        ("Couvert forestier", f"Surface forestière · {dfr['annee_debut']}–"
         f"{dfr['annee_fin']} · % du territoire et km²"),
        ("Inventaire GES 2018", "Émissions par secteur et par gaz, "
         "relues en équivalent CO₂ (PRG 100 ans, AR5)"),
        ("Forêts classées & climat", f"{fo['nb']} forêts classées "
         "(surface, enclavement) et 10 stations météorologiques 2013-2019"),
    ]
    for i, (titre_j, detail) in enumerate(jeux):
        col, lig = i % 2, i // 2
        x = MARGE + col * (LARG / 2 + Inches(.12))
        yy = y + lig * Inches(.92)
        texte(s, x, yy, LARG / 2 - Inches(.2), Inches(.2), titre_j, taille=11,
              gras=True, couleur="encre")
        texte(s, x, yy + Inches(.24), LARG / 2 - Inches(.3), Inches(.5),
              detail, taille=9.5, couleur="sourdine", interligne=1.3)

    encart(s, MARGE, Inches(4.75), LARG, Inches(1.9), "constat",
           "Aucun coût ni budget dans les données : les recommandations "
           "portent sur un ordre de priorité et des cibles physiques, pas sur "
           "un plan de financement.   ·   La cuisson n'est mesurée "
           f"directement qu'en {cb['annees'][0]} et {cb['annees'][1]} : on lit "
           "un contraste, pas une tendance.   ·   La part du recul forestier "
           "imputable au bois-énergie n'est tranchée par aucune source : elle "
           "reste un curseur explicite dans le tableau de bord.",
           titre="Ce que ces données ne permettent pas de dire")

    # ------------------------------------------ 3. diagnostic en 5 chiffres
    s = diapo(prs, "Diagnostic", "Cinq chiffres qui commandent tout le reste", 3)
    larg_t = (LARG - Inches(.48)) / 5
    for i_t, (lab, val, unite, sous, coul) in enumerate(CHIFFRES):
        tuile(s, MARGE + i_t * (larg_t + Inches(.12)), HAUT, larg_t,
              lab, val, unite, sous, coul, h=Inches(1.75))

    encart(s, MARGE, Inches(3.75), LARG, Inches(1.5), "alerte",
           f"Le rythme actuel ne mène pas à 2030, il mène à "
           f"{er['annee_atteinte_tendanciel']:.0f}. L'électrification rurale "
           f"progresse de +{fr(er['rythme_observe'], 2)} point par an depuis "
           f"{er['annee_depart']} ; atteindre 100 % en 2030 en demanderait "
           f"+{fr(er['rythme_requis_2030'], 1)}, soit "
           f"{er['facteur_acceleration']:.0f} fois plus vite. L'extension du "
           f"réseau seule ne produit pas une telle accélération.",
           titre="Le constat qui commande tout le reste")

    # -------------------------------------------------------- 4. la preuve
    s = diapo(prs, "La preuve", "L'accès progresse, la forêt recule, "
              "la cuisson ne bouge pas", 4)
    an_e, v_e = serie_gold("elec_rural", 1998)
    an_c, v_c = serie_gold("cuisson_rural", 1998)
    an_f, v_f = serie_gold("foret_pct", 1998)
    axe = sorted(set(an_e) | set(an_c) | set(an_f))
    courbes(s, MARGE, HAUT, LARG, Inches(2.45), axe,
            [("Accès rural à l'électricité", aligner(axe, an_e, v_e), "energie"),
             ("Accès rural à une cuisson propre",
              aligner(axe, an_c, v_c), "risque")], maxi=40)
    courbes(s, MARGE, Inches(4.24), LARG, Inches(1.86), axe,
            [("Couvert forestier (% du territoire)",
              aligner(axe, an_f, v_f), "foret")], format_val='0.0"%"')
    texte(s, MARGE, Inches(6.16), LARG, Inches(.7),
          f"Deux panneaux, un seul axe du temps : chaque série garde son "
          f"échelle. Le couvert passe de {fr(v_f[0], 1)} % à {fr(v_f[-1], 1)} % du "
          f"territoire pendant que l'accès rural est multiplié par "
          f"{er['valeur']/er['valeur_depart']:.0f} — et que la cuisson propre "
          f"reste à {fr(R['cuisson_rurale']['valeur'], 1)} %. "
          f"L'électricité éclaire les foyers, elle ne remplace pas le bois "
          f"dans les marmites.", taille=10.5, couleur="encre_2",
          interligne=1.35)

    # --------------------------------------------------- 5. électrification
    s = diapo(prs, "Électrification", "Au rythme observé, l'accès universel "
              f"rural tombe en {er['annee_atteinte_tendanciel']:.0f}", 5)
    an_r, v_r = serie_gold("elec_rural", 1998)
    an_u, v_u = serie_gold("elec_urbain", 1998)
    axe2 = sorted(set(an_r) | set(an_u))
    courbes(s, MARGE, HAUT, Inches(7.5), Inches(3.4), axe2,
            [("Urbain", aligner(axe2, an_u, v_u), "urbain"),
             ("Rural", aligner(axe2, an_r, v_r), "energie")], maxi=100)
    xd = MARGE + Inches(7.8)
    ld = LARG - Inches(7.8)
    tuile(s, xd, HAUT, ld, "Écart ville / campagne",
          f"{ec['valeur']:.0f}", "points", f"{ec['urbain']:.0f} % contre "
          f"{ec['rural']:.0f} % en {ec['annee']}", "risque", h=Inches(1.5))
    tuile(s, xd, Inches(3.34), ld, "Coupures subies par les entreprises",
          f"{fr(fi['coupures_mois'], 1)}", "/mois",
          f"{fi['part_entreprises']:.0f} % des entreprises touchées en "
          f"{fi['annee']}, contre {fi['part_entreprises_ref']:.0f} % en "
          f"{fi['annee_ref']}", "risque", h=Inches(1.78))
    encart(s, MARGE, Inches(5.28), LARG, Inches(1.5), "action",
           "Un facteur d'accélération de cet ordre ne s'obtient pas en "
           "prolongeant le réseau, dont le coût au raccordement croît avec la "
           "distance. Il s'obtient par le solaire décentralisé — mini-réseaux "
           "et kits domestiques — déployable sans attendre la ligne moyenne "
           "tension. Second argument : le réseau s'est étendu plus vite qu'il "
           "ne s'est consolidé, la part d'entreprises subissant des coupures "
           "a augmenté.")

    # --------------------------------------------------------- 6. cuisson
    s = diapo(prs, "Cuisson", "Le bois de feu n'est pas une énergie du passé : "
              "sa part augmente", 6)
    an_ec, v_ec = serie_gold("elec_cuisson")
    v0 = [cb["bois"][0], cb["charbon"][0], cb["gpl"][0], v_ec[0]]
    v1 = [cb["bois"][1], cb["charbon"][1], cb["gpl"][1], v_ec[-1]]
    barres(s, MARGE, HAUT, Inches(7.5), Inches(3.5),
           ["Bois de chauffe", "Charbon de bois", "GPL / gaz", "Électricité"],
           [(str(cb["annees"][0]), v0, "neutre"),
            (str(cb["annees"][1]), v1, "risque")], maxi=60)
    tuile(s, MARGE + Inches(7.8), HAUT, LARG - Inches(7.8),
          "Bois brut, entre les deux enquêtes",
          f"+{fr(cb['bois'][1]-cb['bois'][0], 1)}", "points",
          f"de {fr(cb['bois'][0], 1)} % à {fr(cb['bois'][1], 1)} % des ménages, "
          f"quand le charbon perd {fr(cb['charbon'][0]-cb['charbon'][1], 1)} "
          f"points", "risque", h=Inches(1.7))
    tuile(s, MARGE + Inches(7.8), Inches(3.54), LARG - Inches(7.8),
          "Mortalité attribuée à la pollution de l'air",
          f"{sa['mortalite']:.0f}", "",
          f"décès pour 100 000 habitants ({sa['annee_mortalite']}) · PM2,5 à "
          f"{sa['pm25']:.0f} µg/m³, soit × {sa['ratio_oms']:.0f} la ligne "
          f"directrice OMS", "risque", h=Inches(1.7))
    encart(s, MARGE, Inches(5.35), LARG, Inches(1.4), "alerte",
           f"La transition n'a pas commencé, elle a reculé. Le mouvement "
           f"observé est un glissement du charbon vers le bois brut — le "
           f"combustible le plus émetteur de particules à l'intérieur des "
           f"habitations — pendant que le total biomasse reste à "
           f"{cb['biomasse'][1]:.0f} % des ménages.")

    # -------------------------------------------- 7. le piège du renouvelable
    s = diapo(prs, "L'indicateur trompeur",
              "« Renouvelable » ne veut pas dire « propre »", 7)
    an_rn, v_rn = serie_gold("renouv", 1998)
    an_ct, v_ct = serie_gold("cuisson_total", 1998)
    axe3 = sorted(set(an_rn) | set(an_ct))
    courbes(s, MARGE, HAUT, Inches(7.5), Inches(3.5), axe3,
            [("Part « renouvelable » de l'énergie finale",
              aligner(axe3, an_rn, v_rn), "foret"),
             ("Ménages ayant accès à une cuisson propre",
              aligner(axe3, an_ct, v_ct), "risque")], maxi=90)
    texte(s, MARGE + Inches(7.8), HAUT + Inches(.04), LARG - Inches(7.8),
          Inches(3.3),
          [f"En {rp['annee']}, {fr(rp['part_renouvelable'], 1)} % de l'énergie "
           f"finale consommée au Togo est comptée comme renouvelable — un "
           f"taux que peu de pays affichent.",
           f"Au même moment, {fr(rp['part_cuisson_propre'], 1)} % des ménages "
           f"seulement cuisinent proprement.",
           "L'écart entre les deux courbes, c'est la biomasse brûlée dans "
           "les foyers : du bois et du charbon prélevés plus vite qu'ils ne "
           "se régénèrent, brûlés sans évacuation des fumées."],
          taille=11.5, couleur="encre_2", interligne=1.4)
    encart(s, MARGE, Inches(5.35), LARG, Inches(1.4), "action",
           "Conséquence de pilotage : le taux d'énergie renouvelable est un "
           "mauvais indicateur de suivi pour la transition togolaise — il "
           "baissera mécaniquement quand les ménages passeront au GPL ou à "
           "l'électricité, c'est-à-dire quand la situation s'améliorera. "
           "Les indicateurs à suivre sont l'accès à la cuisson propre et la "
           "part des ménages sortis de la biomasse.")

    # ------------------------------------------------------ 8. inventaire
    s = diapo(prs, "Inventaire des émissions", "Le secteur énergie n'est "
              "marginal que si l'on ne regarde que le CO₂", 8)
    gaz = pd.DataFrame(nat["ges_gaz"])
    gaz["co2e"] = gaz["Value"] * gaz["gaz"].map(PRG)
    brut = gaz.groupby("secteur_court")["Value"].sum()
    co2e = gaz.groupby("secteur_court")["co2e"].sum()
    brut_pct = (100 * brut / brut.sum()).sort_values(ascending=False)
    co2e_pct = (100 * co2e / co2e.sum()).reindex(brut_pct.index)
    barres(s, MARGE, HAUT, Inches(7.5), Inches(3.5),
           list(brut_pct.index),
           [("Masse brute (Gg, tel que publié)",
             [round(v, 1) for v in brut_pct], "neutre"),
            ("Équivalent CO₂ (PRG 100 ans)",
             [round(v, 1) for v in co2e_pct], "risque")], maxi=100)
    part_e_co2e = float(co2e_pct.get("Énergie", 0))
    tuile(s, MARGE + Inches(7.8), HAUT, LARG - Inches(7.8),
          "Le secteur énergie, selon l'unité de lecture",
          f"{ges['part_energie']:.0f} → {part_e_co2e:.0f}", "%",
          "des émissions nationales : la même année, le même inventaire, "
          "deux classements différents", "risque", h=Inches(1.7))
    tuile(s, MARGE + Inches(7.8), Inches(3.54), LARG - Inches(7.8),
          "Part de l'énergie dans les gaz les plus réchauffants",
          f"{ges['energie_dans_n2o']:.0f}", "%",
          f"du protoxyde d'azote et {ges['energie_dans_ch4']:.0f} % du "
          f"méthane du pays — deux gaz de combustion incomplète de biomasse",
          "risque", h=Inches(1.7))
    encart(s, MARGE, Inches(5.35), LARG, Inches(1.4), "constat",
           f"L'inventaire publie des masses brutes, où une tonne de méthane "
           f"pèse autant qu'une tonne de CO₂. Or le méthane réchauffe "
           f"{PRG['CH4']} fois plus et le protoxyde d'azote {PRG['N2O']} fois "
           f"plus. « L'énergie ne pèse que {ges['part_energie']:.0f} % » est "
           f"donc une lecture comptable, pas une lecture climatique — et ce "
           f"que cette lecture masque, ce sont les foyers de cuisson.")

    # --------------------------------------------------------- 9. forêts
    s = diapo(prs, "Où agir", f"{fo['nb_robustes']} massifs restent "
              "prioritaires quelle que soit la pondération", 9)
    vul = pd.read_csv(GOLD / "forets_vulnerabilite.csv",
                      usecols=["etab_nom", "region_nom_bdd", "surface_ha",
                               "dist_km", "score"])
    rob = pd.read_csv(GOLD / "forets_robustesse.csv")
    top = (rob[rob["toujours_top10"]]
           .merge(vul, left_on="foret", right_on="etab_nom")
           .sort_values("score", ascending=False))
    noms = [n.replace("Forêt classée ", "").replace("de ", "")
            for n in top["etab_nom"]][::-1]
    barres(s, MARGE, HAUT, Inches(7.6), Inches(4.06), noms,
           [("Indice de vulnérabilité", [round(v, 1) for v in top["score"]][::-1],
             "foret")], legende=False, format_val='0', maxi=100)
    texte(s, MARGE + Inches(7.9), HAUT + Inches(.04), LARG - Inches(7.9), Inches(2.5),
          [f"Sur les {fo['nb']} forêts classées du territoire, "
           f"{fo['nb_robustes']} ne quittent jamais le top 10, quelle que "
           f"soit la pondération testée entre enclavement, surface et stress "
           f"thermique.",
           "Le classement de tête est donc un résultat des données, pas un "
           "choix d'analyste : ce sont de grands massifs enclavés, éloignés "
           "de tout pôle raccordé, donc soumis à la pression du bois-énergie.",
           f"Éloignement médian d'un pôle urbain : {fo['dist_mediane']:.0f} km."],
          taille=11.5, couleur="encre_2", interligne=1.4)
    encart(s, MARGE + Inches(7.9), Inches(4.5), LARG - Inches(7.9),
           Inches(1.28), "action",
           "Ces massifs reçoivent le premier budget de protection. Les forêts "
           "dont l'intervalle de rang est large méritent une instruction "
           "complémentaire avant tout engagement.")

    # ------------------------------------------------ 10. recommandations
    s = diapo(prs, "Recommandations", "Trois leviers, dans cet ordre", 10)
    y = HAUT
    for num, coul, titre_l, quoi, portee, cible, suivi in LEVIERS:
        h = Inches(1.5)
        rectangle(s, MARGE, y, LARG, h, fond="surface", bord="bord")
        texte(s, MARGE + Inches(.22), y + Inches(.3), Inches(.5), Inches(.5),
              num, taille=30, gras=True, couleur=coul, interligne=1)
        texte(s, MARGE + Inches(.75), y + Inches(.2), Inches(3.3), Inches(.3),
              titre_l, taille=14, gras=True, couleur="encre")
        texte(s, MARGE + Inches(.75), y + Inches(.55), Inches(3.4), Inches(.8),
              quoi, taille=9.5, couleur="sourdine", interligne=1.3)
        for i, (etiq, val) in enumerate(
                [("Portée directe", portee), ("Cible 2030", cible),
                 ("Indicateur de suivi", suivi)]):
            x = MARGE + Inches(4.4) + i * Inches(2.6)
            texte(s, x, y + Inches(.22), Inches(2.4), Inches(.18), etiq,
                  taille=8, gras=True, couleur="sourdine", majuscules=True,
                  espacement=.5)
            texte(s, x, y + Inches(.46), Inches(2.42), Inches(.9), val,
                  taille=9.5, couleur="encre", interligne=1.3)
        y += Inches(1.62)

    texte(s, MARGE, Inches(6.6), LARG, Inches(.4),
          "Aucun coût ni budget dans les six jeux de données : ces trois "
          "leviers fixent un ordre de priorité et des cibles physiques, pas "
          "un plan de financement.", taille=9.5, couleur="sourdine",
          interligne=1.3)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    return prs


if __name__ == "__main__":
    # Le compte-rendu reste en ASCII : la console Windows n'est pas en UTF-8
    # et un livrable ne doit pas echouer sur son propre message de fin.
    presentation = construire()
    print(f"[OK] {len(presentation.slides._sldIdLst)} diapositives, "
          f"{OUT.stat().st_size / 1024:.0f} Ko")
    print(f"     {OUT}")
